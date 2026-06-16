"""将 Markdown 格式的幻灯片内容转换为真正的 .pptx 二进制文件"""
import io
import logging
import math
import re
import struct
import zlib
from html import unescape

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from backend.src.utils.slide_schema import parse_markdown_slides


# 颜色方案
DARK_BLUE = RGBColor(0x1A, 0x3C, 0x6E)
ACCENT_BLUE = RGBColor(0x2B, 0x5C, 0x9E)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFC)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)

THEME_PALETTES = {
    "academic_blue": ["#163f8f", "#2f80ed", "#44c2ff", "#f7fbff"],
    "science_green": ["#11695f", "#28b487", "#a7f3d0", "#f6fffb"],
    "warm_case": ["#93491f", "#e86c00", "#ffd166", "#fff8ed"],
    "graphite": ["#17202a", "#566573", "#aeb6bf", "#f7f9fb"],
    "aurora": ["#0f766e", "#22d3ee", "#a78bfa", "#f0fdfa"],
    "coral": ["#9f1239", "#fb7185", "#fbbf24", "#fff1f2"],
    "violet": ["#4c1d95", "#8b5cf6", "#38bdf8", "#f5f3ff"],
    "sunlit": ["#854d0e", "#f59e0b", "#84cc16", "#fffbeb"],
}


def _hex_to_rgb(value: str) -> tuple[int, int, int]:
    value = (value or "#000000").strip().lstrip("#")
    if len(value) != 6:
        return (0, 0, 0)
    return tuple(int(value[i:i + 2], 16) for i in (0, 2, 4))


def _rgb(value: str) -> RGBColor:
    r, g, b = _hex_to_rgb(value)
    return RGBColor(r, g, b)


def _clean_ppt_text(value, limit: int | None = None) -> str:
    text = unescape(str(value or ""))
    text = re.sub(r"<!--[\s\S]*?-->", " ", text)
    text = re.sub(r"</?[^>\n]+>", " ", text)
    text = re.sub(r"<[^>\n]*$", " ", text)
    text = re.sub(r"^\s*(layout|theme|visual)\s*:\s*.*$", " ", text, flags=re.IGNORECASE | re.MULTILINE)
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text).strip()
    if limit and len(text) > limit:
        return text[: max(0, limit - 1)].rstrip() + "..."
    return text


def _fit_font_size(text: str, base: int, dense: int, very_dense: int) -> Pt:
    length = len(_clean_ppt_text(text))
    if length > 120:
        return Pt(very_dense)
    if length > 70:
        return Pt(dense)
    return Pt(base)


def _fit_text_frame(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    return tf


def _split_blocks(items: list[str], size: int) -> list[list[str]]:
    return [items[i:i + size] for i in range(0, len(items), size)]


def _paginate_slides(slides_data: list[dict]) -> list[dict]:
    return slides_data


def _prepare_slide_data(slide_data: dict) -> dict:
    bullets = [
        _clean_ppt_text(item.get("text") if isinstance(item, dict) else item, 110)
        for item in (slide_data.get("bullets") or [])
    ]
    bullets = [item for item in bullets if item]
    blocks = slide_data.get("blocks") if isinstance(slide_data.get("blocks"), list) else []
    if blocks:
        clean_blocks = []
        for block in blocks:
            if isinstance(block, dict):
                text = _clean_ppt_text(block.get("text") or block.get("content"), 110)
                if text:
                    clean_blocks.append({**block, "text": text})
        blocks = clean_blocks
        if not bullets:
            bullets = [block["text"] for block in blocks]

    visual = slide_data.get("visual") if isinstance(slide_data.get("visual"), dict) else {}
    return {
        **slide_data,
        "title": _clean_ppt_text(slide_data.get("title"), 90),
        "text": _clean_ppt_text(slide_data.get("text") or slide_data.get("content"), 900),
        "content": _clean_ppt_text(slide_data.get("content") or slide_data.get("text"), 900),
        "bullets": bullets[:8],
        "notes": _clean_ppt_text(slide_data.get("notes") or slide_data.get("speaker_notes"), 1000),
        "speaker_notes": _clean_ppt_text(slide_data.get("speaker_notes") or slide_data.get("notes"), 1000),
        "blocks": blocks[:8],
        "visual": {
            **visual,
            "query": _clean_ppt_text(visual.get("query"), 90),
            "caption": _clean_ppt_text(visual.get("caption"), 120),
        },
    }


def _palette(slide_data: dict) -> list[str]:
    raw = slide_data.get("palette")
    if isinstance(raw, list) and len(raw) >= 4:
        return raw
    return THEME_PALETTES.get(slide_data.get("theme"), THEME_PALETTES["academic_blue"])


def _chunk(tag: bytes, data: bytes) -> bytes:
    return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)


def _make_visual_png(kind: str, palette: list[str], width: int = 720, height: int = 430) -> bytes:
    c1 = _hex_to_rgb(palette[0])
    c2 = _hex_to_rgb(palette[1])
    c3 = _hex_to_rgb(palette[2])
    rows = []
    kind = kind or "diagram"
    for y in range(height):
        row = bytearray()
        for x in range(width):
            t = (x / max(width - 1, 1) * 0.65) + (y / max(height - 1, 1) * 0.35)
            r = int(c1[0] * (1 - t) + c2[0] * t)
            g = int(c1[1] * (1 - t) + c2[1] * t)
            b = int(c1[2] * (1 - t) + c2[2] * t)
            glow = 0
            if kind == "timeline":
                glow = 90 if abs(y - height * 0.52) < 5 or any((x - width * p) ** 2 + (y - height * 0.52) ** 2 < 28 ** 2 for p in (0.2, 0.4, 0.6, 0.8)) else 0
            elif kind == "comparison":
                glow = 80 if abs(x - width * 0.5) < 5 or (x < width * 0.48 and y > height * 0.22 and y < height * 0.78) or (x > width * 0.52 and y > height * 0.22 and y < height * 0.78) else 0
            elif kind == "formula":
                glow = 95 if abs((y - height * 0.52) - 34 * math.sin(x / 42)) < 4 else 0
            elif kind == "map":
                glow = 85 if ((x - width * 0.34) ** 2 / 19000 + (y - height * 0.48) ** 2 / 7200 < 1) or ((x - width * 0.62) ** 2 / 16000 + (y - height * 0.43) ** 2 / 10000 < 1) else 0
            else:
                glow = 70 if abs((x - width * 0.5) ** 2 / 40000 + (y - height * 0.5) ** 2 / 15000 - 1) < 0.04 or abs(x - y * 1.35) < 4 else 0
            r = min(255, r + int(c3[0] * glow / 255))
            g = min(255, g + int(c3[1] * glow / 255))
            b = min(255, b + int(c3[2] * glow / 255))
            row.extend((r, g, b))
        rows.append(b"\x00" + bytes(row))
    raw = b"".join(rows)
    return b"\x89PNG\r\n\x1a\n" + _chunk(b"IHDR", struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)) + _chunk(b"IDAT", zlib.compress(raw, 9)) + _chunk(b"IEND", b"")


def _parse_slides(markdown: str) -> list[dict]:
    """解析 markdown，返回幻灯片列表 [{title, bullets, notes}]"""
    enriched = parse_markdown_slides(markdown)
    if enriched:
        return [_prepare_slide_data(slide) for slide in enriched]

    raw_slides = re.split(r'\n---\n', markdown.strip())

    slides = []
    first = True
    for block in raw_slides:
        block = block.strip()
        if not block:
            continue

        lines = block.split('\n')
        title = ""
        bullets: list[str] = []
        notes: list[str] = []

        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue

            if stripped.startswith('# '):
                title = stripped[2:].strip()
            elif stripped.startswith('## '):
                title = stripped[3:].strip()
            elif stripped.startswith('> '):
                notes.append(stripped[2:].strip())
            elif stripped.startswith('- ') or stripped.startswith('* '):
                bullets.append(stripped[2:].strip())
            elif re.match(r'^\d+[.)]\s', stripped):
                bullets.append(stripped)
            else:
                bullets.append(stripped)

        if first and not title and bullets:
            title = bullets[0]
            bullets = bullets[1:]

        slides.append(_prepare_slide_data({
            "title": title,
            "bullets": bullets,
            "notes": "\n".join(notes),
        }))
        first = False

    return slides


def _add_dark_bg(slide, prs):
    """深色背景"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()


def _add_light_bg(slide, prs):
    """浅色背景"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()


def _add_theme_bg(slide, prs, slide_data: dict):
    palette = _palette(slide_data)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(palette[3])
    bg.line.fill.background()

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.34))
    band.fill.solid()
    band.fill.fore_color.rgb = _rgb(palette[0])
    band.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.92), prs.slide_width, Inches(0.58))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(palette[1])
    accent.line.fill.background()

    for i, color in enumerate((palette[2], palette[1])):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(9.3 + i * 1.25),
            Inches(0.72 + i * 0.38),
            Inches(2.25 - i * 0.25),
            Inches(2.25 - i * 0.25),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(color)
        shape.fill.transparency = 35 + i * 18
        shape.line.fill.background()


def _add_top_bar(slide, prs):
    """顶部装饰色条"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.12)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()


def _add_title(slide, title_text, left=Inches(0.8), top=Inches(0.4),
               width=Inches(11.7), height=Inches(0.9), font_size=Pt(32)):
    """添加标题"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = _fit_text_frame(box.text_frame)
    p = tf.paragraphs[0]
    clean_title = _clean_ppt_text(title_text, 90)
    p.text = clean_title
    p.font.size = _fit_font_size(clean_title, int(font_size.pt), max(20, int(font_size.pt) - 6), max(17, int(font_size.pt) - 10))
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    return box


def _add_title_line(slide, left=Inches(0.8), top=Inches(1.3), width=Inches(11.7)):
    """标题下分隔线"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()


def _add_bullets(slide, bullets, left=Inches(0.8), top=Inches(1.8),
                 width=Inches(11.7), height=Inches(5.2), max_items=10):
    """添加要点文本，关键术语加粗，层级清晰"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first_item = True
    for i, bullet in enumerate(bullets[:max_items]):
        if first_item:
            p = tf.paragraphs[0]
            first_item = False
        else:
            p = tf.add_paragraph()

        # 彩色序号圆点
        run_dot = p.add_run()
        run_dot.text = "● "
        run_dot.font.size = Pt(10)
        run_dot.font.color.rgb = ACCENT_BLUE
        run_dot.font.vertical_offset = Pt(-2)

        # 尝试按中文冒号分割：加粗前半部分（关键术语）
        parts = bullet.split("：", 1) if "：" in bullet else bullet.split(":", 1)
        if len(parts) == 2:
            run_key = p.add_run()
            run_key.text = parts[0] + "："
            run_key.font.size = Pt(17)
            run_key.font.bold = True
            run_key.font.color.rgb = DARK_BLUE

            run_val = p.add_run()
            run_val.text = parts[1]
            run_val.font.size = Pt(16)
            run_val.font.color.rgb = DARK_GRAY
        else:
            run_val = p.add_run()
            run_val.text = bullet
            run_val.font.size = Pt(16)
            run_val.font.color.rgb = DARK_GRAY

        p.space_after = Pt(6)
        p.space_before = Pt(4)
        p.level = 0
    return box


def _add_bullets(slide, bullets, left=Inches(0.8), top=Inches(1.8),
                 width=Inches(11.7), height=Inches(5.2), max_items=10):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = _fit_text_frame(box.text_frame)
    clean_bullets = [_clean_ppt_text(bullet, 105) for bullet in (bullets or [])]
    clean_bullets = [bullet for bullet in clean_bullets if bullet][:max_items]
    base_size = 14 if len(clean_bullets) <= 5 else 12

    for index, bullet in enumerate(clean_bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        run_dot = p.add_run()
        run_dot.text = "- "
        run_dot.font.size = Pt(max(10, base_size - 2))
        run_dot.font.color.rgb = ACCENT_BLUE

        parts = bullet.split("：", 1) if "：" in bullet else bullet.split(":", 1)
        if len(parts) == 2 and parts[0].strip() and parts[1].strip():
            run_key = p.add_run()
            run_key.text = parts[0].strip() + ": "
            run_key.font.size = _fit_font_size(bullet, base_size + 1, base_size, max(9, base_size - 3))
            run_key.font.bold = True
            run_key.font.color.rgb = DARK_BLUE

            run_val = p.add_run()
            run_val.text = parts[1].strip()
            run_val.font.size = _fit_font_size(bullet, base_size, max(10, base_size - 1), max(8, base_size - 3))
            run_val.font.color.rgb = DARK_GRAY
        else:
            run_val = p.add_run()
            run_val.text = bullet
            run_val.font.size = _fit_font_size(bullet, base_size, max(10, base_size - 1), max(8, base_size - 3))
            run_val.font.color.rgb = DARK_GRAY

        p.space_after = Pt(3 if len(clean_bullets) >= 6 else 5)
        p.space_before = Pt(2)
        p.level = 0
    return box


def _add_bullets(slide, bullets, left=Inches(0.8), top=Inches(1.8),
                 width=Inches(11.7), height=Inches(5.2), max_items=10):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = _fit_text_frame(box.text_frame)
    clean_bullets = [_clean_ppt_text(bullet, 92) for bullet in (bullets or [])]
    clean_bullets = [bullet for bullet in clean_bullets if bullet][:max_items]
    if not clean_bullets:
        return box

    if len(clean_bullets) >= 6:
        font_size = 10
        after = 2
    elif len(clean_bullets) >= 4:
        font_size = 11
        after = 3
    else:
        font_size = 12
        after = 4

    for index, bullet in enumerate(clean_bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"- {bullet}"
        p.font.size = _fit_font_size(bullet, font_size, max(9, font_size - 1), 8)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(after)
        p.space_before = Pt(0)
        p.level = 0
    return box


def _build_title_slide(prs, slide_data: dict):
    """标题页 — 深色背景 + 居中主标题 + 副标题"""
    layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(layout)
    _add_dark_bg(slide, prs)

    # 装饰元素 - 右上角小色块
    decor = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(10.5), Inches(0), Inches(2.8), Inches(0.06)
    )
    decor.fill.solid()
    decor.fill.fore_color.rgb = ACCENT_ORANGE
    decor.line.fill.background()

    # 标题
    box = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.5))
    tf = _fit_text_frame(box.text_frame)
    p = tf.paragraphs[0]
    p.text = slide_data["title"] or "课件"
    p.text = _clean_ppt_text(p.text, 90)
    p.font.size = _fit_font_size(p.text, 44, 36, 30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 副标题或第一要点
    subtitle_text = ""
    if slide_data["bullets"]:
        subtitle_text = slide_data["bullets"][0]
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.5))
        stf = _fit_text_frame(sub_box.text_frame)
        sp = stf.paragraphs[0]
        sp.text = subtitle_text
        sp.text = _clean_ppt_text(sp.text, 130)
        sp.font.size = _fit_font_size(sp.text, 22, 18, 15)
        sp.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        sp.alignment = PP_ALIGN.CENTER

        for extra in slide_data["bullets"][1:3]:
            sp2 = stf.add_paragraph()
            sp2.text = extra
            sp2.text = _clean_ppt_text(sp2.text, 120)
            sp2.font.size = _fit_font_size(sp2.text, 18, 15, 12)
            sp2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
            sp2.alignment = PP_ALIGN.CENTER

    _add_notes(slide, slide_data)


def _build_content_slide(prs, slide_data: dict, index: int, total: int):
    """纯内容页：标题 + 要点 + 页码"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _add_theme_bg(slide, prs, slide_data)
    _add_top_bar(slide, prs)
    _add_visual_thumb(slide, slide_data)
    _add_title(slide, slide_data["title"] or f"第 {index} 页")
    _add_title_line(slide)
    _add_bullets(slide, slide_data["bullets"])
    _add_page_number(slide, index, total)
    _add_notes(slide, slide_data)


def _add_visual_panel(slide, x, y, w, h, caption: str = "", accent=ACCENT_BLUE, slide_data: dict | None = None):
    palette = _palette(slide_data or {})
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = _rgb(palette[3])
    panel.line.color.rgb = _rgb(palette[1])

    visual = (slide_data or {}).get("visual") or {}
    png = _make_visual_png(visual.get("type", "diagram"), palette)
    slide.shapes.add_picture(io.BytesIO(png), x + Inches(0.25), y + Inches(0.25), width=w - Inches(0.5), height=h - Inches(1.25))

    box = slide.shapes.add_textbox(x + Inches(0.55), y + h - Inches(0.92), w - Inches(1.1), Inches(0.58))
    tf = _fit_text_frame(box.text_frame)
    p = tf.paragraphs[0]
    p.text = _clean_ppt_text(caption, 90)
    p.font.size = _fit_font_size(p.text, 13, 11, 10)
    p.font.color.rgb = MEDIUM_GRAY
    p.alignment = PP_ALIGN.CENTER


def _add_visual_thumb(slide, slide_data: dict, left=Inches(10.25), top=Inches(5.2), width=Inches(2.45), height=Inches(1.45)):
    palette = _palette(slide_data)
    visual = slide_data.get("visual") or {}
    png = _make_visual_png(visual.get("type", "diagram"), palette, width=520, height=310)
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left - Inches(0.08), top - Inches(0.08), width + Inches(0.16), height + Inches(0.16))
    frame.fill.solid()
    frame.fill.fore_color.rgb = _rgb(palette[3])
    frame.line.color.rgb = _rgb(palette[1])
    slide.shapes.add_picture(io.BytesIO(png), left, top, width=width, height=height)


def _build_concept_visual_slide(prs, slide_data: dict, index: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_theme_bg(slide, prs, slide_data)
    _add_top_bar(slide, prs)
    _add_title(slide, slide_data["title"] or f"Slide {index}", width=Inches(6.2), font_size=Pt(30))
    _add_title_line(slide, width=Inches(6.0))
    _add_bullets(slide, slide_data.get("bullets", []), left=Inches(0.8), top=Inches(1.65), width=Inches(6.1), height=Inches(5.1), max_items=5)
    visual = slide_data.get("visual") or {}
    _add_visual_panel(slide, Inches(7.3), Inches(1.35), Inches(5.25), Inches(4.9), visual.get("caption") or visual.get("query") or "", slide_data=slide_data)
    _add_page_number(slide, index, total)
    _add_notes(slide, slide_data)


def _build_process_slide(prs, slide_data: dict, index: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_theme_bg(slide, prs, slide_data)
    _add_top_bar(slide, prs)
    _add_visual_thumb(slide, slide_data)
    _add_title(slide, slide_data["title"] or f"Slide {index}", font_size=Pt(30))
    steps = (slide_data.get("bullets") or [])[:4]
    start_x = Inches(1.0)
    y = Inches(2.35)
    card_w = Inches(2.55)
    gap = Inches(0.28)
    for i, step in enumerate(steps):
        x = start_x + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, Inches(3.0))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xC7, 0xD8, 0xE8)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.86), y - Inches(0.35), Inches(0.82), Inches(0.82))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT_ORANGE if i % 2 else ACCENT_BLUE
        badge.line.fill.background()
        n = slide.shapes.add_textbox(x + Inches(0.86), y - Inches(0.22), Inches(0.82), Inches(0.4))
        p = n.text_frame.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.48), card_w - Inches(0.36), Inches(2.15))
        tf = _fit_text_frame(box.text_frame)
        p = tf.paragraphs[0]
        p.text = _clean_ppt_text(step, 80)
        p.font.size = _fit_font_size(p.text, 11, 9, 8)
        p.font.color.rgb = DARK_GRAY
    _add_page_number(slide, index, total)
    _add_notes(slide, slide_data)


def _build_comparison_slide(prs, slide_data: dict, index: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_theme_bg(slide, prs, slide_data)
    _add_top_bar(slide, prs)
    _add_visual_thumb(slide, slide_data)
    _add_title(slide, slide_data["title"] or f"Slide {index}", font_size=Pt(30))
    items = (slide_data.get("bullets") or [])[:6]
    left_items = items[::2]
    right_items = items[1::2] or items[3:]
    for col, (x, color, heading) in enumerate(((Inches(0.9), ACCENT_BLUE, "A"), (Inches(6.85), ACCENT_ORANGE, "B"))):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.65), Inches(5.55), Inches(4.95))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xC7, 0xD8, 0xE8)
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.65), Inches(5.55), Inches(0.55))
        head.fill.solid()
        head.fill.fore_color.rgb = color
        head.line.fill.background()
        hb = slide.shapes.add_textbox(x + Inches(0.24), Inches(1.74), Inches(5.05), Inches(0.35))
        hp = hb.text_frame.paragraphs[0]
        hp.text = heading
        hp.font.size = Pt(15)
        hp.font.bold = True
        hp.font.color.rgb = WHITE
        _add_bullets(slide, left_items if col == 0 else right_items, left=x + Inches(0.35), top=Inches(2.45), width=Inches(4.9), height=Inches(3.6), max_items=3)
    _add_page_number(slide, index, total)
    _add_notes(slide, slide_data)


def _build_formula_slide(prs, slide_data: dict, index: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_theme_bg(slide, prs, slide_data)
    _add_top_bar(slide, prs)
    _add_visual_thumb(slide, slide_data)
    _add_title(slide, slide_data["title"] or f"Slide {index}", font_size=Pt(30))
    bullets = slide_data.get("bullets") or []
    formula = next((b for b in bullets if "=" in b or "$" in b), bullets[0] if bullets else "")
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.75), Inches(10.9), Inches(1.65))
    panel.fill.solid()
    panel.fill.fore_color.rgb = DARK_BLUE
    panel.line.fill.background()
    box = slide.shapes.add_textbox(Inches(1.55), Inches(2.08), Inches(10.2), Inches(0.95))
    _fit_text_frame(box.text_frame)
    p = box.text_frame.paragraphs[0]
    p.text = _clean_ppt_text(formula.replace("$", ""), 140)
    p.font.size = _fit_font_size(p.text, 26, 21, 17)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    rest = [b for b in bullets if b != formula]
    _add_bullets(slide, rest, left=Inches(1.2), top=Inches(3.8), width=Inches(10.8), height=Inches(2.55), max_items=4)
    _add_page_number(slide, index, total)
    _add_notes(slide, slide_data)


def _build_cards_slide(prs, slide_data: dict, index: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_theme_bg(slide, prs, slide_data)
    _add_top_bar(slide, prs)
    _add_visual_thumb(slide, slide_data)
    _add_title(slide, slide_data["title"] or f"Slide {index}", font_size=Pt(30))
    items = (slide_data.get("bullets") or [])[:6]
    for i, item in enumerate(items):
        row = i // 3
        col = i % 3
        x = Inches(0.85 + col * 4.15)
        y = Inches(1.65 + row * 2.25)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.75), Inches(1.78))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xC7, 0xD8, 0xE8)
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.12), Inches(1.78))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = ACCENT_BLUE if i % 2 == 0 else ACCENT_ORANGE
        stripe.line.fill.background()
        box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.22), Inches(3.15), Inches(1.26))
        tf = _fit_text_frame(box.text_frame)
        p = tf.paragraphs[0]
        p.text = _clean_ppt_text(item, 82)
        p.font.size = _fit_font_size(p.text, 10, 9, 8)
        p.font.color.rgb = DARK_GRAY
    _add_page_number(slide, index, total)
    _add_notes(slide, slide_data)


def _add_page_number(slide, index: int, total: int):
    """右下角页码"""
    box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{index} / {total}"
    p.font.size = Pt(11)
    p.font.color.rgb = MEDIUM_GRAY
    p.alignment = PP_ALIGN.RIGHT


def _add_notes(slide, slide_data: dict):
    """添加备注"""
    if slide_data["notes"]:
        try:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = _clean_ppt_text(slide_data["notes"], 1000)
        except Exception:
            logging.getLogger("pptx").warning("添加备注失败 slide_title=%s", slide_data.get("title", ""))


def markdown_to_pptx(markdown: str) -> bytes:
    """将 markdown 幻灯片内容转换为 .pptx 二进制数据"""
    slides_data = _paginate_slides(_parse_slides(markdown))
    if not slides_data:
        return b""

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = len(slides_data)
    for i, slide_data in enumerate(slides_data):
        if i == 0:
            _build_title_slide(prs, slide_data)
        elif slide_data.get("layout") == "concept_visual":
            _build_concept_visual_slide(prs, slide_data, i, total)
        elif slide_data.get("layout") == "process_steps":
            _build_process_slide(prs, slide_data, i, total)
        elif slide_data.get("layout") == "comparison":
            _build_comparison_slide(prs, slide_data, i, total)
        elif slide_data.get("layout") == "formula_focus":
            _build_formula_slide(prs, slide_data, i, total)
        elif slide_data.get("layout") == "content_cards":
            _build_cards_slide(prs, slide_data, i, total)
        else:
            _build_content_slide(prs, slide_data, i, total)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
