"""测试完整管线：topic → 讯飞 PPTX"""
import asyncio
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent.parent))
DESKTOP = Path("C:/Users/35808/Desktop")


async def main():
    from backend.src.service.ppt_service import PptService

    query = (
        "请生成一份关于「STM32F103 单片机核心架构与原理」的详细 PPT。\n"
        "目标受众：大学计算机专业学生，已掌握基础编程知识。\n"
        "教学要求：中等难度，侧重硬件原理与架构理解。\n"
        "要求：内容丰富、每页要点充足，适合课堂教学。"
        "请自动配图，包含架构图、流程图等 visual 元素。"
        "至少 10 页以上。"
    )

    print("[INFO] 调用讯飞智文 API 生成 PPT...")
    try:
        local_path = await PptService.generate(query, is_figure=True)
        src = Path(local_path)
        dst = DESKTOP / "STM32_Final.pptx"
        dst.write_bytes(src.read_bytes())
        size_kb = src.stat().st_size / 1024
        print(f"[OK] 已保存到桌面: {dst} ({size_kb:.1f} KB)")

        # 验证
        from pptx import Presentation as PPTXReader
        from pptx.shapes.picture import Picture
        prs = PPTXReader(str(dst))
        chars, imgs = 0, 0
        for i, slide in enumerate(prs.slides):
            pics = [s for s in slide.shapes if isinstance(s, Picture)]
            texts = [p.text for s in slide.shapes if s.has_text_frame for p in s.text_frame.paragraphs if p.text.strip()]
            sc = sum(len(t) for t in texts)
            chars += sc
            imgs += len(pics)
            title = texts[0][:50] if texts else "(empty)"
            print(f"  Slide {i}: {len(pics)} img, {sc}c - {title}")
        print(f"  [TOTAL] {len(prs.slides)} slides, {imgs} images, {chars} chars")
    except Exception as e:
        print(f"[ERROR] {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    asyncio.run(main())
