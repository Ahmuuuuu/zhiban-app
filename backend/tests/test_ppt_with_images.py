"""测试：生成带配图的 PPT 并保存到桌面"""
import asyncio
import sys
import re
import os
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent.parent))

DESKTOP = Path(os.environ.get("USERPROFILE", "C:/Users/35808")) / "Desktop"
TOPIC = "CPU 缓存架构与工作原理"


async def main():
    # 1. 初始化数据库
    from backend.src.utils.database import init_db, close_db
    await init_db()

    # 2. 找一个用户
    from backend.src.models.usermodel import User
    user = await User.first()
    if not user:
        print("[ERROR] 数据库中没有用户，请先注册")
        return
    user_id = user.id
    print(f"[OK] 使用用户: {user.username} (id={user_id})")

    # 3. 构造 PPT prompt
    from backend.src.utils.prompt_loader import load_prompt, fill_prompt

    template = load_prompt("resource/ppt")
    system_prompt = fill_prompt(
        template,
        portrait_context="大学计算机专业学生，已掌握基础编程知识",
        learning_guidance="中等难度，侧重硬件原理与架构理解",
        kb_context="",
        feedback="",
        topic=TOPIC,
    )
    print(f"[INFO] Prompt 已构建 ({len(system_prompt)} chars)")

    # 4. 调用 LLM 生成 Markdown
    from backend.src.ai_core.llm_config import llm
    from langchain_core.messages import HumanMessage, SystemMessage

    print("[WAIT] 等待 LLM 生成 PPT 内容...")
    resp = await llm.ainvoke([
        SystemMessage(content=system_prompt),
        HumanMessage(content=f"请生成关于「{TOPIC}」的 PPT 内容"),
    ])
    markdown = resp.content
    print(f"[OK] LLM 返回 {len(markdown)} chars")

    # 5. 保存原始 markdown 到桌面（调试用）
    (DESKTOP / "_ppt_raw.md").write_text(markdown, encoding="utf-8")
    print(f"[INFO] 原始 Markdown 已保存到桌面 _ppt_raw.md")

    # 6. 检查配图标记
    img_markers = list(re.finditer(r'<!--\s*image:\s*(.+?)\s*-->', markdown, re.IGNORECASE))
    print(f"[INFO] 检测到 {len(img_markers)} 个配图标记")
    for m in img_markers:
        print(f"   - {m.group(1).strip()[:60]}")

    # 7. 生成图片
    from backend.src.service.image_service import ImageService

    # 先解析 slides，确定哪些 slide index 需要图
    from backend.src.utils.pptx_generator import _parse_slides
    slides_data = _parse_slides(markdown)

    image_paths: dict[int, str] = {}
    static_dir = Path(__file__).parent.parent / "static"

    # 遍历 slides，有 image_desc 的就生图
    for i, slide in enumerate(slides_data):
        desc = slide.get("image_desc", "")
        if not desc:
            continue
        print(f"[WAIT] 生成第 {i} 页配图: [{desc[:40]}]...")
        try:
            images = await ImageService.generate(desc, str(user_id), "4:3", 1)
            if images and images[0].get("url"):
                url = images[0]["url"]  # /static/images/xxx.jpg
                local_path = static_dir / "images" / url.rsplit("/", 1)[-1]
                if local_path.exists():
                    image_paths[i] = str(local_path)
                    print(f"  [OK] 图片已保存: {local_path}")
                else:
                    print(f"  [WARN] URL 返回了但本地文件不存在: {url}")
            else:
                print(f"  [WARN] 图片生成失败: {images}")
        except Exception as e:
            print(f"  [ERROR] 图片生成异常: {e}")

    # 8. 生成 PPTX
    from backend.src.utils.pptx_generator import markdown_to_pptx

    print(f"[INFO] 生成 PPTX ({len(image_paths)} 张配图)...")
    pptx_bytes = markdown_to_pptx(markdown)
    if not pptx_bytes:
        print("[ERROR] PPTX 生成为空")
        return

    save_path = DESKTOP / "CPU_Cache_PPT.pptx"
    save_path.write_bytes(pptx_bytes)
    print(f"[OK] PPTX 已保存: {save_path}")
    print(f"     {len(slides_data)} 页, {len(image_paths)} 张配图")

    # 9. 验证 PPTX 中的图片
    from pptx import Presentation as PPTXReader
    prs = PPTXReader(str(save_path))
    total_chars = 0
    total_images = 0
    for i, slide in enumerate(prs.slides):
        from pptx.shapes.picture import Picture
        pics = [s for s in slide.shapes if isinstance(s, Picture)]
        texts = []
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
        title = texts[0] if texts else "(no text)"
        slide_chars = sum(len(t) for t in texts)
        total_chars += slide_chars
        total_images += len(pics)
        status = "OK" if len(pics) > 0 else "MISSING"
        print(f"  [VERIFY] Slide {i}: {len(pics)} images ({status}), {slide_chars} chars - {title[:40]}")
    print(f"  [SUMMARY] {total_images} images, {total_chars} total chars across {len(prs.slides)} slides")

    await close_db()


if __name__ == "__main__":
    asyncio.run(main())
